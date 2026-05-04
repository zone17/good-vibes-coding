"""
Transform GoodVibes-Coding-Business-Plan-v2.pptx -> v3 (full pass).

Strategy: surgical text replacement keyed by (slide_index, shape_name).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

SRC = Path("/Users/fp/Desktop/good-vibes-coding/docs/strategy/2026-04-30-business-plan/deck/GoodVibes-Coding-Business-Plan-v2.pptx")
DST = Path("/Users/fp/Desktop/good-vibes-coding/docs/strategy/2026-04-30-business-plan/deck/GoodVibes-Coding-Business-Plan-v3.pptx")
NEW_TOTAL = 52


def replace_text_keep_format(shape, new_text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        tf.text = new_text
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = new_text
        for run in list(first_para.runs[1:]):
            run._r.getparent().remove(run._r)
    else:
        run = first_para.add_run()
        run.text = new_text
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)


def replace_multiline(shape, new_text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        tf.text = new_text
        return
    lines = new_text.split("\n")
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = lines[0]
        for run in list(first_para.runs[1:]):
            run._r.getparent().remove(run._r)
    else:
        run = first_para.add_run()
        run.text = lines[0]
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line


# =================================================================
# Edits keyed by (slide_index 1-based, shape_name) -> new_text
# =================================================================
EDITS = {
    # Slide 1 — Cover
    (1, "Text 1"): "Business plan, v3.",
    (1, "Text 2"): "Cold-acquisition pivot. Path A only — productized 2-person, no SaaS spinoff. Two geographies — Palm Beach + Poconos. Three lower-ticket starter offerings ($750 / $1,500 / $1,800). Workshops + GBP + podcast as the cold-acquisition stack. Seven open questions that unblock the first thirty days.",

    # Slide 2 — Plan in one page
    (2, "Text 3"): "Stay 2-person and productize (Path A only). Three cold-conversion starters that ladder to recurring retainers. WikiFold stays free as proof of capability. Cold acquisition driven by local SEO + workshops + owned podcast. No spinoff. No outside capital.",
    (2, "Text 4"): "What v3 changed",
    (2, "Text 5"): "Cold acquisition replaces the warm-referral primary motion. Path C (WikiFold spinoff) is dropped. Two explicit geographies — Palm Beach + Poconos. Lower starter prices match the cold-conversion ceiling. Workshop circuit becomes channel #2.",
    (2, "Text 7"): "Y3 base $612K · $288K combined draws · 30% margin · ~33-hour weeks. Y1 base ~$510K (workshop sprint volume layers on top of recurring revenue).",
    (2, "Text 9"): "Path A · Productized 2-person",
    (2, "Text 10"): "$700K – 1.4M",
    (2, "Text 11"): "Path C · Y3+ optionality only",
    (2, "Text 12"): "$14M – 55M only if WikiFold productizes later",
    (2, "Text 13"): "Path B · Boutique hire-up (rejected)",
    (2, "Text 14"): "$1.8M – 3M",
    (2, "Text 15"): "Path D · Methodology license (deferred Y3+)",
    (2, "Text 16"): "$3M – 18M",
    (2, "Text 17"): "$20K MRR by month 9–10 base case. Workshop conversion + sprint→retainer rate are the two levers.",
    (2, "Text 18"): f"02 / {NEW_TOTAL} · EXECUTIVE SUMMARY",

    # Slide 4 — Thesis sources
    (4, "Text 8"): "Source: 01-thesis-validation.md, 03-market-sizing.md. Goldman Sachs 10KSB, Feb 2026; SBA Office of Advocacy 2025.",

    # Slide 9 — Unique abilities (soften Mel's role)
    (9, "Text 1"): "Operator instinct + curious-listener voice. Poconos field credibility.",

    # Slide 10 — First three Whos REORDERED
    (10, "Text 0"): "FIRST THREE WHOS",
    (10, "Text 1"): "The cold-motion forces a resequence of who we hire first.",
    (10, "Text 4"): "Content + Ops VA — month 2–3",
    (10, "Text 5"): "~$2,500–3,000/mo. Pulls forward in v3. Owns workshop logistics + podcast distribution + direct-mail follow-up + lead tracking. Removes the multi-market admin load before it stalls the founders.",
    (10, "Text 6"): "Fractional podcast editor — month 6–9 (conditional)",
    (10, "Text 7"): "~$2,500/mo. Conditional on podcast pulling its weight. Skip entirely if M6 podcast-attributable leads are zero.",
    (10, "Text 8"): "Local Poconos field rep — month 12 (non-negotiable)",
    (10, "Text 9"): "$1,500/mo, part-time. Cold email scales to zero in the Poconos market. Without local presence, the pipeline collapses.",

    # Slide 11 — Part III opener
    (11, "Text 2"): "Path A only.\nConfirmed.",
    (11, "Text 3"): "Stay 2-person and productize. WikiFold stays free as proof. Two geographies. Cold-acquisition stack. Paths B, C, D, E, F all evaluated and rejected as primary plan.",

    # Slide 12 — The recommended path (drop Path C column)
    (12, "Text 1"): "Path A — operator-led, 2-person, productized.",
    (12, "Text 2"): "Productized 2-person · months 0–18+",
    (12, "Text 3"): "Three cold-conversion starters → retainer ladder. Workshops + local SEO + owned podcast as the channel stack. Two geographies. Sub-$2K starter price band. Up-tier sprints unlocked once trust is earned in a market.",
    (12, "Text 4"): "What we're NOT doing · v3",
    (12, "Text 5"): "No SaaS spinoff. WikiFold stays free as proof of capability. No multi-tier subscriptions. No outside capital. No equity / rev-share. No relocation. No cold outbound to Fortune 5000. Never lead with 'AI consulting'.",
    (12, "Text 6"): "(deferred)",
    (12, "Text 7"): "Path C — Y3+ optionality only",
    (12, "Text 8"): "WikiFold productization revisited only if customer demand signals at M24+. Walk-away $14–55M only if that productization happens — not in the primary plan.",
    (12, "Text 9"): "The podcast sits across consulting and free-tools — distributing methodology IP, building audience, serving as the launch surface for any future productization decision.",

    # Slide 16 — Pricing menu (full rewrite of all 9 SKUs)
    (16, "Text 0"): "PRICING MENU · 9 SKUS",
    (16, "Text 1"): "The full menu, front to back.",
    (16, "Text 2"): "THREE COLD-CONVERSION STARTERS",
    (16, "Text 3"): "AI Visibility Tune-Up",
    (16, "Text 4"): "Replaces Yext, BrightLocal · 5 days",
    (16, "Text 5"): "$750",
    (16, "Text 6"): "+ $250/mo Watch",
    (16, "Text 7"): "Menu + Reviews Lite",
    (16, "Text 8"): "Replaces Bloom + Birdeye + OpenMenu · 1 wk",
    (16, "Text 9"): "$1,500",
    (16, "Text 10"): "+ $300/mo Hospitality Lite",
    (16, "Text 11"): "AI Front-Desk",
    (16, "Text 12"): "Replaces Goodcall, Numa, Rosie · 2 wks",
    (16, "Text 13"): "$1,800",
    (16, "Text 14"): "+ $200/mo built-in",
    (16, "Text 15"): "BACKUPS · MONTHS 2–4",
    (16, "Text 16"): "Vacation Rental Direct-Booking",
    (16, "Text 17"): "Pocs-led · beats Airbnb take + Hostfully",
    (16, "Text 18"): "$2,400",
    (16, "Text 19"): "+ $400/mo",
    (16, "Text 20"): "SOP Lite",
    (16, "Text 21"): "$2,200",
    (16, "Text 22"): "+ $300/mo wiki keep-alive",
    (16, "Text 23"): "Restaurant Full Stack · BUNDLE",
    (16, "Text 24"): "Replaces Yext + Birdeye + OpenMenu + Goodcall · M3+ anchor",
    (16, "Text 25"): "$3,500",
    (16, "Text 26"): "+ $500/mo",
    (16, "Text 27"): "Hospitality Premium · post-trust",
    (16, "Text 28"): "B&B Stack Consolidator, Wedding Inquiry, Outfitter Booking",
    (16, "Text 29"): "$3.5–7.5K",
    (16, "Text 30"): "+ $600–850/mo",
    (16, "Text 31"): "RETAINERS · THE LADDER",
    (16, "Text 32"): "Watch · GBP active monitoring",
    (16, "Text 33"): "AI search visibility · monthly local SEO updates",
    (16, "Text 34"): "$250/mo",
    (16, "Text 35"): "—",
    (16, "Text 36"): "Hospitality Lite",
    (16, "Text 37"): "Menu + reviews + allergen ongoing",
    (16, "Text 38"): "$300/mo",
    (16, "Text 39"): "—",
    (16, "Text 40"): f"16 / {NEW_TOTAL} · PART IV — REVENUE MODEL",

    # Slide 17 — Three starter offerings
    (17, "Text 0"): "THE THREE STARTER OFFERINGS",
    (17, "Text 1"): "Three cold-conversion starters. One retainer ladder.",
    (17, "Text 2"): "Anti-SaaS",
    (17, "Text 3"): "AI Visibility Tune-Up",
    (17, "Text 4"): "Claim and optimize Google Business Profile + AI search schema. Cold-pitch frame: 'You're paying Yext or BrightLocal $99–$300/mo for the same thing — we set it up once, you own it.'",
    (17, "Text 5"): "$750 flat · 5 days",
    (17, "Text 6"): "Ladders to $250/mo Watch retainer.",
    (17, "Text 7"): "Hospitality wedge",
    (17, "Text 8"): "Menu + Reviews Lite",
    (17, "Text 9"): "Menu single-source-of-truth + AI review responses + allergen tagging. Replaces Bloom + Birdeye + OpenMenu (~$300+/mo). CA SB-68 (Jul 1, 2026) is the FL regulatory tailwind.",
    (17, "Text 10"): "$1,500 flat · 1 week",
    (17, "Text 11"): "Ladders to $300/mo Hospitality Lite retainer.",
    (17, "Text 12"): "Service-business wedge",
    (17, "Text 13"): "AI Front-Desk",
    (17, "Text 14"): "AI phone receptionist + call routing for HVAC, dental, salons, B&Bs. Replaces Goodcall / Numa / Rosie ($49–$79/mo). Built-in retainer.",
    (17, "Text 15"): "$1,800 flat · 2 weeks",
    (17, "Text 16"): "Built-in $200/mo retainer for ongoing call routing.",

    # Slide 18 — Anti-patterns refused
    (18, "Text 1"): "What we evaluated and chose not to build (v3).",
    (18, "Text 4"): "$7,500 site refresh (v2 starter) · standalone",
    (18, "Text 5"): "Mispriced for cold motion · cold-conversion ceiling is ~$3K not $7.5K",
    (18, "Text 6"): "Customer-service chatbots",
    (18, "Text 7"): "Qualtrics: 4× failure rate of other AI tasks · 88% of customers prefer humans · 'fire your staff with AI' anti-pattern",
    (18, "Text 8"): "Mailchimp → Klaviyo migrations",
    (18, "Text 9"): "E-commerce-tilted · saturated done-for-you category · off-vertical for hospitality + services",
    (18, "Text 10"): "Custom booking-system rebuilds",
    (18, "Text 11"): "Vertical SaaS does it cheaper · $40K–$250K customs are out of starter band",
    (18, "Text 12"): "Standalone bookkeeping automation",
    (18, "Text 13"): "Regulated · QuickBooks AI handles 85–95% already · build a partner channel instead",
    (18, "Text 14"): "Hourly $200/hr 'AI consulting' on the homepage",
    (18, "Text 15"): "The 10x trap · reserved as $450/hr internal-only existing-client rate",

    # Slide 19 — Part V opener
    (19, "Text 2"): "Cold acquisition wins.\nWorkshops + GBP + podcast lead.",
    (19, "Text 3"): "The warm-referral motion is upside, not plan. The owned podcast is the most exciting asset, not the most productive lead source. Workshops + local SEO + direct mail (Poconos) carry the pipeline.",

    # Slide 20 — Channel ranking (full rewrite — re-rank for cold motion)
    (20, "Text 1"): "Where leads actually come from in year one — cold motion.",
    (20, "Text 7"): "Local SEO + GBP (two SABs)",
    (20, "Text 8"): "Matt",
    (20, "Text 9"): "3–6 months",
    (20, "Text 11"): "Workshop circuit (libraries, Chambers, coworking)",
    (20, "Text 12"): "Both",
    (20, "Text 13"): "4–8 weeks",
    (20, "Text 15"): "Owned podcast + newsletter recap",
    (20, "Text 16"): "Both (Mel-led audiograms)",
    (20, "Text 17"): "8–16 weeks",
    (20, "Text 19"): "Chamber / BNI / Rotary engagement",
    (20, "Text 20"): "Both",
    (20, "Text 21"): "6–12 weeks",
    (20, "Text 23"): "Direct mail (Poconos only)",
    (20, "Text 24"): "Matt designs, Mel signs",
    (20, "Text 25"): "4–8 weeks",
    (20, "Text 27"): "Cold direct outreach (email + LinkedIn + Loom)",
    (20, "Text 28"): "Both",
    (20, "Text 29"): "4–8 weeks",
    (20, "Text 31"): "Sponsored events / festivals / fairs",
    (20, "Text 32"): "Both",
    (20, "Text 33"): "Seasonal",
    (20, "Text 35"): "Geo-targeted Google / Meta ads",
    (20, "Text 36"): "Matt",
    (20, "Text 37"): "2–4 weeks",

    # Slide 22 — Honest 12-month math
    (22, "Text 10"): "Honest framing: workshops + local SEO + cold outreach carry the bulk of pipeline. Podcast adds ~25–35% of attributed Y1 revenue at 1,200 listeners/ep base case. Newsletter recap is the actual conversion lever, not the audio.",

    # Slide 36 — Three scenarios (update Y1)
    (36, "Text 1"): "Three scenarios. Y1 base ~$510K with workshop sprint volume layered on.",
    (36, "Text 15"): "~$510K",
    (36, "Text 16"): "~$720K",

    # Slide 38 — Sensitivity
    (38, "Text 1"): "Top 4 levers (v3): client wins, sprint price, workshop conversion, churn.",
    (38, "Text 11"): "Workshop conversion · NEW IN v3",
    (38, "Text 12"): "Sprints/workshop × sprint→retainer rate. The dominant lever in cold motion.",

    # Slide 39 — The one number
    (39, "Text 1"): "The one number: $20K MRR by month 9–10.",
    (39, "Text 2"): "Workshop conversion is the metric to watch weekly. Below 3 sprints/workshop or below 50% sprint→retainer by M3, audit the offer before scaling outreach.",

    # Slide 41 — First 90 days (rewrite per v3 §6)
    (41, "Text 1"): "First 30 days, week by week.",
    (41, "Text 2"): "Week 1",
    (41, "Text 3"): "Site live with 3 cold-conversion starters ($750 / $1,500 / $1,800) + new tagline. Two GBPs claimed (PB + Pocs SAB). 23 cold emails sent (Mel 13 hospitality, Matt 8 eng-peer + 10 cold). First 2 workshops booked.",
    (41, "Text 4"): "Week 2",
    (41, "Text 5"): "10 more cold emails per geography. First 2 booked Calendly calls. Direct mail (Pocs) postcard list finalized. Podcast at 4×/week + Friday short. Newsletter recap goes live.",
    (41, "Text 6"): "Weeks 3–4",
    (41, "Text 7"): "First sprint sold (likely AI Visibility Tune-Up at $750). First sample-build Loom recorded. Direct mail wave 1 ships (200 postcards). 4 booked calls cumulative.",
    (41, "Text 8"): "Weeks 5–8",
    (41, "Text 9"): "First sprint shipped. Retainer offered as yes/no follow-on. Second sprint sold. First workshop runs in PB (target 20+ attendees, 3+ individual consults).",
    (41, "Text 10"): "Weeks 9–12",
    (41, "Text 11"): "First Operations Retainer signed (~$250–$3K MRR). 4 workshops cadence stable (2/geography/month). 8–10 sprints sold MTD. Restaurant Full Stack bundle launches as anchor offer. First case study published.",
    (41, "Text 12"): "Day-30 exit metrics: 40+ named conversations · 6+ Calendly calls · 2 sprints sold · 1 shipped · 1 retainer offered · 1 workshop run · 200 postcards sent · 2 GBPs claimed.",

    # Slide 42 — 12-month quarterly view
    (42, "Text 1"): "12-month quarterly view (v3 cold motion).",
    (42, "Text 9"): "Bookkeeper · Content + Ops VA",
    (42, "Text 13"): "—",
    (42, "Text 17"): "Fractional podcast editor (conditional) · Ops / PM coordinator",
    (42, "Text 21"): "Local Poconos field rep ($1,500/mo)",
    (42, "Text 22"): "End of Year 1 base case ceiling: ~$510K revenue · 12–15 clients · 3 fractional hires running · workshop circuit producing pipeline · podcast at 1.2–1.5K listeners/ep.",

    # Slide 43 — Decision triggers (simplify to M3 / M6 / M9)
    (43, "Text 1"): "Decision triggers: M3 / M6 / M9.",
    (43, "Text 6"): "3",
    (43, "Text 7"): "4+ workshops run, 3+ sprints/workshop, 3+ retainers signed?",
    (43, "Text 8"): "Commit cadence",
    (43, "Text 9"): "Audit offer",
    (43, "Text 10"): "3",
    (43, "Text 11"): "Direct mail (Pocs) at 0.5%+ response rate?",
    (43, "Text 12"): "Continue + scale to 500-postcard waves",
    (43, "Text 13"): "Pause Pocs cold; lean on workshops only",
    (43, "Text 14"): "6",
    (43, "Text 15"): "$10K+ MRR running?",
    (43, "Text 16"): "Continue current motion",
    (43, "Text 17"): "Drop Pocs; concentrate Palm Beach",
    (43, "Text 18"): "6",
    (43, "Text 19"): "Podcast-attributable leads ≥1/month?",
    (43, "Text 20"): "Hire podcast editor",
    (43, "Text 21"): "Skip editor; downscale podcast",
    (43, "Text 22"): "6",
    (43, "Text 23"): "Founders' weekly hours <60 each?",
    (43, "Text 24"): "Sustainable",
    (43, "Text 25"): "Hire 2nd VA; renegotiate cadence",
    # Second column (M9 / sponsorship / Path A check)
    (43, "Text 30"): "9",
    (43, "Text 31"): "$20K+ MRR running?",
    (43, "Text 32"): "Land junior implementer + Pocs field rep",
    (43, "Text 33"): "Hold; tighten conversion before adding payroll",
    (43, "Text 34"): "9",
    (43, "Text 35"): "Path A still right path?",
    (43, "Text 36"): "Continue",
    (43, "Text 37"): "Re-evaluate",
    (43, "Text 38"): "18",
    (43, "Text 39"): "First sponsorship offer arrived?",
    (43, "Text 40"): "Refuse before M18 regardless",
    (43, "Text 41"): "—",
    (43, "Text 42"): "—",
    (43, "Text 43"): "—",
    (43, "Text 44"): "—",
    (43, "Text 45"): "—",
    (43, "Text 46"): "—",
    (43, "Text 47"): "—",
    (43, "Text 48"): "—",
    (43, "Text 49"): "—",
    (43, "Text 50"): "—",
    (43, "Text 51"): "—",
    (43, "Text 52"): "—",
    (43, "Text 53"): "—",

    # Slide 46 — Open questions (v3 — 7 questions)
    (46, "Text 0"): "OPEN QUESTIONS",
    (46, "Text 1"): "Seven open questions to unblock the first 30 days.",
    (46, "Text 2"): "Pricing commit",
    (46, "Text 3"): "$750 / $1,500 / $1,800 starters live-publishable? Including the $250 / $300 / built-in $200 retainer tiers?",
    (46, "Text 4"): "Workshop venue commits",
    (46, "Text 5"): "Which library / Chamber / coworking for the first PB workshop? Which for the first Pocs workshop?",
    (46, "Text 6"): "First Poconos residency week",
    (46, "Text 7"): "When? Recommendation: late February 2027 or earlier. Earlier is better for the M3 decision point.",
    (46, "Text 8"): "Direct mail list commit",
    (46, "Text 9"): "200-postcard first wave + letter follow-up. Budget ~$300 first wave / ~$1,000 first month.",
    (46, "Text 10"): "Two Google Business Profiles",
    (46, "Text 11"): "Service-Area Business setup OK with both founders' addresses? Matt FL primary, Mel address for Pocs lookup.",
    (46, "Text 12"): "VA hiring commit",
    (46, "Text 13"): "$2,500–$3,000/mo at M2–3. Confirm budget tolerance.",
    (46, "Text 14"): "Tagline lock",
    (46, "Text 15"): "'Stop renting your business from five SaaS companies. Own it once.' OK as the cross-market headline?",
    (46, "Text 16"): "—",
    (46, "Text 17"): "—",
    (46, "Text 18"): "—",
    (46, "Text 19"): "—",
    (46, "Text 20"): "—",
    (46, "Text 21"): "—",

    # Slide 47 — Open questions · podcast & content commitments
    (47, "Text 0"): "OPEN QUESTIONS · PODCAST & FORMAT",
    (47, "Text 1"): "Podcast and format commitments.",
    (47, "Text 2"): "Cadence commit",
    (47, "Text 3"): "4×/week (Mon–Thu) + Friday short, or hold daily? `15-podcast-strategy.md` recommends 4×/wk to recover ~6 hr/wk.",
    (47, "Text 4"): "Editorial framework lock",
    (47, "Text 5"): "'The Augmentation Lens' or 'The Three-Question Test' or other? Without a named frame, the editorial position stays copyable.",
    (47, "Text 6"): "Sign-off ritual lock",
    (47, "Text 7"): "Is 'Good vibes only' OK as the closing line? Lock for 8 weeks for habit-formation.",
    (47, "Text 8"): "Podcast editor budget commit",
    (47, "Text 9"): "$2,500/mo from M6–9 — conditional on podcast-attributable leads ≥1 by M6. Skip if below.",
    (47, "Text 10"): "First-five hospitality customers",
    (47, "Text 11"): "Actual names, not just neighborhoods. Mel's network list from week 1.",
    (47, "Text 12"): "CA SB-68 angle yes/no",
    (47, "Text 13"): "Does the Menu + Reviews Lite sprint lean into the July 2026 allergen-disclosure law as the headline pitch?",
    (47, "Text 14"): "Multi-market rotation cadence",
    (47, "Text 15"): "2-week PB / 2-week Pocs rotation, 60/40 PB lean through M6 — confirm willingness.",

    # Slide 48 — What changed vs v1 (now: vs v2)
    (48, "Text 0"): "WHAT CHANGED VS V2",
    (48, "Text 1"): "What changed vs v2.",

    # Slide 49 — Closing posture
    (49, "Text 1"): "Path A only. Cold acquisition. Two geographies.",
    (49, "Text 2"): "Stop renting your business from five SaaS companies. Own it once.",
    (49, "Text 3"): "Three cold-conversion starters at $750 / $1,500 / $1,800. Workshops + GBP + podcast as the channel stack. Path C deferred. Mel's network is upside, not plan. The seven open questions in the prior slide unblock the first 30 days.",
    (49, "Text 4"): "The next move is the founders'.",
}


def update_footer_counters(prs):
    """Update '/49' or '/ 49' in footers to NEW_TOTAL."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    if "/ 49" in t:
                        run.text = t.replace("/ 49", f"/ {NEW_TOTAL}")
                    elif "/49" in t:
                        run.text = t.replace("/49", f"/{NEW_TOTAL}")


def add_text_only_slide(prs, label, title, body, source_line=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    left = Inches(1.0)
    width = Inches(18.0)
    # Title
    tx = slide.shapes.add_textbox(left, Inches(1.5), width, Inches(2.0))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    if tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].font.size = Pt(54)
        tf.paragraphs[0].runs[0].font.bold = True
    # Body
    body_box = slide.shapes.add_textbox(left, Inches(4.0), width, Inches(5.5))
    btf = body_box.text_frame
    btf.word_wrap = True
    body_lines = body.split("\n")
    btf.paragraphs[0].text = body_lines[0]
    if btf.paragraphs[0].runs:
        btf.paragraphs[0].runs[0].font.size = Pt(28)
    for line in body_lines[1:]:
        p = btf.add_paragraph()
        p.text = line
        if p.runs:
            p.runs[0].font.size = Pt(28)
    # Footer
    foot_box = slide.shapes.add_textbox(left, Inches(10.5), width, Inches(0.5))
    foot_tf = foot_box.text_frame
    foot_tf.paragraphs[0].text = label
    if foot_tf.paragraphs[0].runs:
        foot_tf.paragraphs[0].runs[0].font.size = Pt(14)
    if source_line:
        src_box = slide.shapes.add_textbox(left, Inches(9.7), width, Inches(0.5))
        src_tf = src_box.text_frame
        src_tf.paragraphs[0].text = source_line
        if src_tf.paragraphs[0].runs:
            src_tf.paragraphs[0].runs[0].font.size = Pt(14)
            src_tf.paragraphs[0].runs[0].font.italic = True


def main():
    prs = Presentation(str(SRC))
    print(f"Loaded {SRC.name} — {len(prs.slides)} slides")

    applied = 0
    skipped = 0
    skipped_keys = []
    for (slide_idx, shape_name), new_text in EDITS.items():
        if slide_idx > len(prs.slides):
            skipped_keys.append((slide_idx, shape_name, "OOR"))
            skipped += 1
            continue
        slide = prs.slides[slide_idx - 1]
        match = next((s for s in slide.shapes if s.name == shape_name), None)
        if match is None:
            skipped_keys.append((slide_idx, shape_name, "no shape"))
            skipped += 1
            continue
        if "\n" in new_text:
            replace_multiline(match, new_text)
        else:
            replace_text_keep_format(match, new_text)
        applied += 1

    print(f"  Applied: {applied}, Skipped: {skipped}")
    if skipped_keys:
        for k in skipped_keys[:10]:
            print(f"    skipped: {k}")

    # Add NEW slides
    add_text_only_slide(
        prs,
        label=f"50 / {NEW_TOTAL} · TAGLINE",
        title="Stop renting your business from five SaaS companies.",
        body="Own it once.\n\nAnti-SaaS is the brand spine. Every offering replaces a named subscription with code the customer owns.\n\nYext · BrightLocal · Bloom · Birdeye · OpenMenu · Goodcall · Numa · Hostfully — every one of them is a target.",
    )
    add_text_only_slide(
        prs,
        label=f"51 / {NEW_TOTAL} · TWO GEOGRAPHIES",
        title="Two geographies. Differentiated pricing.",
        body=(
            "Palm Beach County, FL  ·  ~45–50K employer firms across 9 cities\n"
            "  Hospitality + home services + professional services + med spas + marine\n"
            "  Pricing: standard tier ($750 / $1,500 / $1,800 starters · $250–500/mo retainer)\n"
            "  Buying window: May–July\n"
            "\n"
            "Narrowsburg / Poconos region (Sullivan NY + Wayne/Pike/Monroe/Carbon PA)\n"
            "  ~422,650 residents · $8.2B/yr visitor spend\n"
            "  B&Bs + outfitters + restaurants + makers + farms + service trades\n"
            "  Pricing: 30% lower · $750–$1,200/mo retainer floor · $300/mo total-software 'ouch' line\n"
            "  Buying window: Feb 1 – Apr 15 off-season\n"
            "  3 residency weeks/yr · $1,500/mo field rep at M12 — non-negotiable"
        ),
        source_line="Source: 19-south-florida-market.md, 20-poconos-upstate-market.md",
    )
    add_text_only_slide(
        prs,
        label=f"52 / {NEW_TOTAL} · $20K MRR RAMP",
        title="$20K MRR by month 9–10. Honest math.",
        body=(
            "Base case (workshops + cold + paid layered):\n"
            "  M3 ~$2K  ·  M5 ~$8K  ·  M7 ~$13K  ·  M9 ~$18K  ·  M10 ~$21K  ✓\n"
            "\n"
            "Upside case (cold paid converts AI Front-Desk at 4–5%):\n"
            "  $20K MRR by M7–8\n"
            "\n"
            "Conservative case (workshop attendance <15 or sprint→retainer <40%):\n"
            "  $20K MRR by M11–13\n"
            "\n"
            "Year-1 base case revenue: ~$510K (workshop sprint volume layers on top of recurring)."
        ),
        source_line="Source: 21-product-portfolio-options.md, 22-cold-acquisition-gtm.md",
    )

    update_footer_counters(prs)
    prs.save(str(DST))
    print(f"\nSaved → {DST}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
